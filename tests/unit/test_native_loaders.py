# Copyright 2024-2026 Firefly Software Foundation
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Unit tests for the native per-format SourceLoaders."""

from __future__ import annotations

import io

from flycanon.core.services.ingestion.loaders import (
    CsvLoader,
    JsonLoader,
    MarkdownLoader,
    OdfLoader,
    PptxLoader,
    RtfLoader,
    TextLoader,
    XlsxLoader,
    XmlLoader,
    default_registry,
)
from flycanon.interfaces.enums import SourceKind


def _body(doc) -> str:
    return "\n".join(s.body for s in doc.sections) if doc.sections else (doc.raw_text or "")


# --- delimited / structured text -------------------------------------------


def test_csv_loader_renders_rows() -> None:
    doc = CsvLoader().load("name,age\nAda,36\nGrace,40\n")
    body = _body(doc)
    assert "name\tage" in body
    assert "Ada\t36" in body


def test_tsv_loader_uses_tab_delimiter() -> None:
    loader = CsvLoader(kind=SourceKind.tsv, delimiter="\t")
    assert loader.kind == SourceKind.tsv
    doc = loader.load("a\tb\n1\t2\n")
    assert "a\tb" in _body(doc)


def test_json_loader_pretty_prints() -> None:
    doc = JsonLoader().load('{"b": 2, "a": [1, 2]}')
    body = _body(doc)
    assert '"b": 2' in body
    assert '"a"' in body


def test_json_loader_falls_back_to_raw_on_invalid() -> None:
    doc = JsonLoader().load("not valid json {")
    assert "not valid json" in _body(doc)


def test_xml_loader_extracts_text() -> None:
    doc = XmlLoader().load("<root><a>hello</a><b>world</b></root>")
    body = _body(doc)
    assert "hello" in body
    assert "world" in body


def test_rtf_loader_strips_control_words() -> None:
    doc = RtfLoader().load(r"{\rtf1\ansi\deff0 Hello RTF World}")
    body = _body(doc)
    assert "Hello RTF World" in body
    assert "rtf1" not in body


# --- binary office formats (built in-memory with the same libs) ------------


def test_xlsx_loader_one_section_per_sheet() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Numbers"
    ws.append(["x", "y"])
    ws.append([1, 2])
    buf = io.BytesIO()
    wb.save(buf)

    doc = XlsxLoader().load(buf.getvalue())
    assert doc.sections
    assert doc.sections[0].path == ["Numbers"]
    assert "x\ty" in doc.sections[0].body


def test_pptx_loader_one_section_per_slide() -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello Slide One"
    buf = io.BytesIO()
    prs.save(buf)

    doc = PptxLoader().load(buf.getvalue())
    assert any("Hello Slide One" in s.body for s in doc.sections)
    assert doc.sections[0].path == ["Slide 1"]


def test_odt_loader_extracts_body_text() -> None:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    document = OpenDocumentText()
    document.text.addElement(P(text="Hello ODT body"))
    buf = io.BytesIO()
    document.save(buf)

    doc = OdfLoader(kind=SourceKind.odt).load(buf.getvalue())
    assert "Hello ODT body" in _body(doc)


# --- registry wiring --------------------------------------------------------


def test_default_registry_has_native_loaders_and_text_fallback() -> None:
    registry = default_registry()
    assert isinstance(registry.get(SourceKind.xlsx), XlsxLoader)
    assert isinstance(registry.get(SourceKind.pptx), PptxLoader)
    assert isinstance(registry.get(SourceKind.csv), CsvLoader)
    assert isinstance(registry.get(SourceKind.tsv), CsvLoader)
    assert isinstance(registry.get(SourceKind.json_), JsonLoader)
    assert isinstance(registry.get(SourceKind.xml), XmlLoader)
    assert isinstance(registry.get(SourceKind.rtf), RtfLoader)
    assert isinstance(registry.get(SourceKind.odt), OdfLoader)
    assert isinstance(registry.get(SourceKind.ods), OdfLoader)
    assert isinstance(registry.get(SourceKind.odp), OdfLoader)
    # Multi-artifact archives arrive pre-merged as Markdown.
    assert isinstance(registry.get(SourceKind.archive), MarkdownLoader)
    # Unrecognised kinds degrade to the plain-text fallback.
    assert isinstance(registry.get(SourceKind.unknown), TextLoader)
