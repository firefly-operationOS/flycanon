# Copyright 2026 Firefly Software Solutions Inc
"""Source loaders -- normalised bytes in, :class:`LoadedDocument` out.

Each loader implements :class:`SourceLoader` and is responsible for
producing a normalised representation of one source: a list of
:class:`Section` objects carrying the heading hierarchy, the verbatim
text under each section, and a small bit of positional metadata
(``page``, ``order``).

Every supported format has a dedicated, dependency-light loader with no
Microsoft MarkItDown dependency: DOCX (python-docx), PDF (PyMuPDF +
Tesseract OCR), HTML (BeautifulSoup), Markdown / Text, XLSX (openpyxl),
PPTX (python-pptx), CSV / TSV (stdlib), JSON (stdlib), XML (lxml), RTF
(striprtf), ODT / ODS / ODP (odfpy), images (Tesseract OCR) and
transcripts (WebVTT / SRT). Multi-artifact intakes (archives, emails)
arrive pre-merged as Markdown and are loaded by :class:`MarkdownLoader`.
The registry's fallback is a plain UTF-8 :class:`TextLoader`, so an
unrecognised payload degrades to text rather than failing.
"""

from __future__ import annotations

import contextlib
import io
import logging
import re
from dataclasses import dataclass, field
from typing import Protocol

from bs4 import BeautifulSoup
from PIL import Image as _PILImage

from flycanon.core.services.ingestion.errors import CorruptSource
from flycanon.interfaces.enums import SourceKind

logger = logging.getLogger(__name__)


@dataclass(slots=True)
class Section:
    """A heading + body pair recovered by the loader.

    ``path`` is the breadcrumb of headings leading to this section
    (e.g. ``["Scope", "In scope"]``). ``order`` keeps the loader's
    emission order so chunks remain stable across reruns.
    """

    path: list[str]
    body: str
    order: int = 0
    page: int | None = None


@dataclass(slots=True)
class LoadedDocument:
    """Normalised view of a source the chunker can consume."""

    sections: list[Section] = field(default_factory=list)
    page_count: int | None = None
    language: str | None = None
    title: str | None = None
    raw_text: str | None = None


class SourceLoader(Protocol):
    """Bytes (or pre-decoded text) -> :class:`LoadedDocument`."""

    kind: SourceKind

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument: ...


# ---------------------------------------------------------------------------
# Implementations
# ---------------------------------------------------------------------------


class DocxLoader:
    """High-fidelity DOCX loader on top of ``python-docx``.

    Walks the paragraph stream, tracks ``Heading1`` ... ``Heading9``
    styles, and emits one section per heading. Body paragraphs are
    coalesced under the most recent heading.
    """

    kind = SourceKind.docx

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            from docx import Document
        except ImportError as exc:
            raise CorruptSource("docx", f"python-docx is not installed: {exc}") from exc

        try:
            if isinstance(content, str):
                content = content.encode("utf-8")
            document = Document(io.BytesIO(content))
        except Exception as exc:
            raise CorruptSource("docx", str(exc)) from exc

        sections: list[Section] = []
        heading_stack: list[str] = []
        current_path: list[str] = []
        current_body: list[str] = []
        order = 0

        def _flush() -> None:
            nonlocal order
            body = "\n".join(line for line in current_body if line.strip())
            if not body:
                return
            sections.append(Section(path=list(current_path), body=body, order=order))
            order += 1

        title: str | None = None
        for para in document.paragraphs:
            style = (para.style.name if para.style else "") or ""
            text = (para.text or "").strip()
            if not text:
                continue
            heading_match = re.match(r"Heading\s*(\d+)", style)
            if style == "Title":
                title = title or text
                continue
            if heading_match:
                _flush()
                current_body = []
                level = int(heading_match.group(1))
                while len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(text)
                current_path = list(heading_stack)
                continue
            current_body.append(text)
        _flush()

        if not sections:
            return LoadedDocument(title=title)
        return LoadedDocument(sections=sections, title=title)


class PdfLoader:
    """PDF loader supporting both PDF kinds with a two-phase strategy.

    flycanon ingests two flavours of PDF:

    * **Full Digital Text PDFs** -- born-digital documents whose text
      content is stored as encoded glyphs in the page stream
      (Word / LibreOffice / LaTeX exports, browser "Save as PDF",
      PDFs produced by reporting pipelines). Text extraction is
      lossless and effectively free.

    * **PDF-Images (scanned PDFs)** -- pages that carry only a
      raster image of the original document (scanned contracts, fax
      output, photos of receipts, mobile-camera captures). The page
      stream has no encoded text; reading the words requires Optical
      Character Recognition.

    Phase 1 -- text-layer extraction. PyMuPDF (``pymupdf``, aka
    ``fitz``) returns the document's encoded text stream per page in
    microseconds without rendering. For born-digital PDFs that's all
    we need.

    Phase 2 -- per-page OCR fallback. Pages whose extracted text is
    shorter than :attr:`_MIN_CHARS_PER_PAGE` are flagged as
    image-only. For each flagged page we rasterise via PyMuPDF at
    :attr:`_OCR_DPI` and pipe the PNG bytes into Tesseract
    (``pytesseract.image_to_string``). The OCR language is composed
    as ``eng+spa`` by default and overridable via
    ``FLYCANON_OCR_LANG``. The engine is selectable via
    ``FLYCANON_PDF_OCR_ENGINE`` (``tesseract`` default; ``docling``
    after installing the ``docling`` extra for layout-aware OCR with
    native multi-column / table handling). The two phases compose:
    a hybrid PDF (text on some pages, scans on others) gets text
    extraction where available and OCR where needed.

    pypdf is intentionally NOT used here -- it is reserved for
    :class:`PdfGuard`'s lightweight encryption / corruption
    pre-flight. MarkItDown is also not on the PDF path; it remains
    the last-resort fallback for exotic formats only.
    """

    kind = SourceKind.pdf

    # Pages with fewer than this many extracted characters trigger
    # the OCR render-and-recognise path.
    _MIN_CHARS_PER_PAGE = 16
    # Tesseract DPI for rendering -- 200 is a good speed/accuracy
    # trade-off for typical scanned business documents.
    _OCR_DPI = 200

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        if isinstance(content, str):
            content = content.encode("utf-8")

        try:
            import pymupdf
        except ImportError as exc:  # pragma: no cover -- runtime dep guard
            raise RuntimeError("pymupdf is required for PDF ingestion (replaces pypdf/MarkItDown)") from exc

        try:
            doc = pymupdf.open(stream=content, filetype="pdf")
        except Exception as exc:
            raise CorruptSource("pdf", str(exc)) from exc

        if doc.needs_pass:
            doc.close()
            raise CorruptSource("pdf", "encrypted PDFs are not supported")

        page_count = doc.page_count

        # Phase 1: PyMuPDF text-layer extraction. ``get_text()`` with
        # no mode returns the document's encoded text stream in
        # reading order; for born-digital PDFs that's all we need.
        text_per_page: list[str] = []
        try:
            for page_index in range(page_count):
                try:
                    raw = doc[page_index].get_text("text") or ""
                    text = raw.strip() if isinstance(raw, str) else ""
                except Exception:
                    text = ""
                text_per_page.append(text)
        finally:
            doc.close()

        # Phase 2: identify image-only pages and OCR them.
        needs_ocr = [idx for idx, t in enumerate(text_per_page) if len(t) < self._MIN_CHARS_PER_PAGE]
        if needs_ocr:
            ocr_results = self._ocr_pages(content, needs_ocr)
            for page_idx, ocr_text in ocr_results.items():
                if ocr_text.strip():
                    text_per_page[page_idx] = ocr_text.strip()

        sections: list[Section] = []
        for page_idx, text in enumerate(text_per_page):
            if not text:
                continue
            sections.append(
                Section(
                    path=[f"Page {page_idx + 1}"],
                    body=text,
                    order=page_idx,
                    page=page_idx + 1,
                )
            )
        return LoadedDocument(sections=sections, page_count=page_count)

    def _ocr_pages(self, pdf_bytes: bytes, page_indices: list[int]) -> dict[int, str]:
        """Render ``page_indices`` and OCR them.

        Engine selection: ``FLYCANON_PDF_OCR_ENGINE`` (default
        ``tesseract``; set to ``docling`` after installing the
        ``docling`` extra to get layout-aware OCR with native
        multi-column / table handling). Both engines rasterise via
        PyMuPDF at the configured DPI.
        """
        import os as _os

        engine = _os.environ.get("FLYCANON_PDF_OCR_ENGINE", "tesseract").strip().lower()
        if engine == "docling":
            results = self._ocr_pages_docling(pdf_bytes, page_indices)
            if results:
                return results
            logger.warning("docling OCR returned no pages; falling back to tesseract")
        return self._ocr_pages_tesseract(pdf_bytes, page_indices)

    def _ocr_pages_tesseract(self, pdf_bytes: bytes, page_indices: list[int]) -> dict[int, str]:
        try:
            import os as _os

            import pymupdf
            import pytesseract
        except ImportError as exc:
            logger.warning("PDF OCR fallback unavailable: %s", exc)
            return {}

        lang = _os.environ.get("FLYCANON_OCR_LANG", "eng+spa")
        results: dict[int, str] = {}
        try:
            doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        except Exception as exc:
            logger.warning("PyMuPDF could not open PDF for OCR: %s", exc)
            return {}

        try:
            for page_idx in page_indices:
                try:
                    page = doc[page_idx]
                    pix = page.get_pixmap(dpi=self._OCR_DPI)
                    png_bytes = pix.tobytes("png")
                    ocr_text = pytesseract.image_to_string(
                        _PILImage.open(io.BytesIO(png_bytes)),
                        lang=lang,
                    )
                    results[page_idx] = ocr_text
                except Exception as exc:
                    logger.warning("PDF OCR failed for page %d: %s", page_idx + 1, exc)
                    results[page_idx] = ""
        finally:
            with contextlib.suppress(Exception):
                doc.close()
        return results

    def _ocr_pages_docling(self, pdf_bytes: bytes, page_indices: list[int]) -> dict[int, str]:
        """Layout-aware OCR via IBM Docling. Requires the ``docling``
        extra: ``uv pip install flycanon[docling]``.
        """
        try:
            from docling.datamodel.base_models import DocumentStream  # type: ignore[import-not-found]
            from docling.document_converter import DocumentConverter  # type: ignore[import-not-found]
        except ImportError as exc:
            logger.warning(
                "docling not installed; install the ``docling`` extra "
                "to enable FLYCANON_PDF_OCR_ENGINE=docling (%s)",
                exc,
            )
            return {}

        try:
            converter = DocumentConverter()
            stream = DocumentStream(name="ingested.pdf", stream=io.BytesIO(pdf_bytes))
            doc_result = converter.convert(stream)
            # Docling represents the document as a tree of items each
            # carrying a page index. We project them down to per-page
            # plain text in reading order.
            buckets: dict[int, list[str]] = {}
            for item in getattr(doc_result.document, "iterate_items", lambda: [])():
                page = getattr(item, "prov", None)
                page_idx = 0
                if page and isinstance(page, list) and page:
                    page_idx = max(0, int(getattr(page[0], "page_no", 1)) - 1)
                text = getattr(item, "text", None) or ""
                if text.strip():
                    buckets.setdefault(page_idx, []).append(text)
            return {idx: "\n".join(parts).strip() for idx, parts in buckets.items() if idx in page_indices}
        except Exception as exc:  # noqa: BLE001 -- fall back to tesseract
            logger.warning("docling OCR failed: %s", exc)
            return {}


class HtmlLoader:
    """BeautifulSoup-backed HTML loader.

    Heading tags (``h1`` ... ``h6``) drive section boundaries; their
    text becomes the section path. Paragraph + list-item nodes are
    coalesced under the most recent heading.
    """

    kind = SourceKind.html

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        try:
            soup = BeautifulSoup(content, "lxml")
        except Exception:
            soup = BeautifulSoup(content, "html.parser")

        title = soup.title.string.strip() if soup.title and soup.title.string else None

        sections: list[Section] = []
        heading_stack: list[str] = []
        current_path: list[str] = []
        current_body: list[str] = []
        order = 0

        def _flush() -> None:
            nonlocal order
            body = "\n".join(line for line in current_body if line.strip())
            if not body:
                return
            sections.append(Section(path=list(current_path), body=body, order=order))
            order += 1

        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li"]):
            text = tag.get_text(" ", strip=True)
            if not text:
                continue
            name = getattr(tag, "name", None) or ""
            if name.startswith("h") and name[1:].isdigit():
                _flush()
                current_body = []
                level = int(name[1:])
                while len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(text)
                current_path = list(heading_stack)
                continue
            current_body.append(text)
        _flush()

        return LoadedDocument(sections=sections, title=title)


class MarkdownLoader:
    """Plain-text walk over Markdown.

    Headings are detected by line prefix (``# `` ... ``###### ``); the
    leading body becomes a section with an empty ``path``.
    """

    kind = SourceKind.markdown

    _HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        return _markdown_to_document(content)


class TextLoader:
    """Plain-text loader -- the whole content lands as one section."""

    kind = SourceKind.text

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        body = content.strip()
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class XlsxLoader:
    """XLSX loader on top of ``openpyxl`` -- one section per worksheet.

    Each non-empty worksheet becomes a section whose body is the cell
    grid rendered as tab-separated rows (blank rows dropped). Formulae
    are read as their cached values (``data_only=True``).
    """

    kind = SourceKind.xlsx

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:  # pragma: no cover -- runtime dep guard
            raise CorruptSource("xlsx", f"openpyxl is not installed: {exc}") from exc
        if isinstance(content, str):
            content = content.encode("utf-8")
        try:
            workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        except Exception as exc:
            raise CorruptSource("xlsx", str(exc)) from exc

        sections: list[Section] = []
        order = 0
        try:
            for sheet in workbook.worksheets:
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    cells = ["" if value is None else str(value) for value in row]
                    if any(cell.strip() for cell in cells):
                        rows.append("\t".join(cells).rstrip())
                body = "\n".join(rows).strip()
                if body:
                    sections.append(Section(path=[sheet.title], body=body, order=order))
                    order += 1
        finally:
            with contextlib.suppress(Exception):
                workbook.close()

        if not sections:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=sections)


class PptxLoader:
    """PPTX loader on top of ``python-pptx`` -- one section per slide.

    Walks every shape's text frame (and table cells) per slide; the
    coalesced text becomes the slide's section body.
    """

    kind = SourceKind.pptx

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover -- runtime dep guard
            raise CorruptSource("pptx", f"python-pptx is not installed: {exc}") from exc
        if isinstance(content, str):
            content = content.encode("utf-8")
        try:
            presentation = Presentation(io.BytesIO(content))
        except Exception as exc:
            raise CorruptSource("pptx", str(exc)) from exc

        sections: list[Section] = []
        for index, slide in enumerate(presentation.slides):
            lines: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    for table_row in shape.table.rows:
                        cells = [cell.text.strip() for cell in table_row.cells]
                        if any(cells):
                            lines.append("\t".join(cells))
            body = "\n".join(lines).strip()
            if body:
                sections.append(Section(path=[f"Slide {index + 1}"], body=body, order=index))
        if not sections:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=sections)


class CsvLoader:
    """Delimited-text loader (CSV / TSV) on top of the stdlib ``csv`` module.

    Serves both :data:`SourceKind.csv` and :data:`SourceKind.tsv`; pass
    ``kind`` + ``delimiter`` to bind the instance. The whole grid lands as
    one section rendered as tab-separated rows.
    """

    def __init__(self, *, kind: SourceKind = SourceKind.csv, delimiter: str = ",") -> None:
        self.kind = kind
        self._delimiter = delimiter

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        import csv

        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        try:
            rows = list(csv.reader(io.StringIO(content), delimiter=self._delimiter))
        except csv.Error as exc:
            raise CorruptSource(self.kind.value, str(exc)) from exc
        lines = ["\t".join(cell.strip() for cell in row) for row in rows if any(cell.strip() for cell in row)]
        body = "\n".join(lines).strip()
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class JsonLoader:
    """JSON loader -- pretty-prints valid JSON; falls back to raw text.

    Pretty-printing (sorted-free, indented) gives the chunker stable,
    line-oriented structure to split on without inventing a schema.
    """

    kind = SourceKind.json_

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        import json

        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        try:
            body = json.dumps(json.loads(content), indent=2, ensure_ascii=False)
        except (ValueError, TypeError):
            body = content.strip()
        body = body.strip()
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class XmlLoader:
    """XML loader -- extracts visible text via ``lxml`` (recovering parser).

    Falls back to the raw decoded text when the document cannot be parsed
    so malformed XML still ingests as plain text.
    """

    kind = SourceKind.xml

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        raw = content.decode("utf-8", errors="replace") if isinstance(content, bytes) else content
        body = raw.strip()
        try:
            from lxml import etree

            root = etree.fromstring(raw.encode("utf-8"), parser=etree.XMLParser(recover=True))
            if root is not None:
                parts = [text.strip() for text in root.itertext() if text and text.strip()]
                if parts:
                    body = "\n".join(parts)
        except Exception:  # noqa: BLE001 -- malformed XML degrades to raw text
            pass
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class RtfLoader:
    """RTF loader on top of ``striprtf`` (pure-Python, no LibreOffice)."""

    kind = SourceKind.rtf

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError as exc:  # pragma: no cover -- runtime dep guard
            raise CorruptSource("rtf", f"striprtf is not installed: {exc}") from exc
        raw = content.decode("utf-8", errors="replace") if isinstance(content, bytes) else content
        try:
            body = (rtf_to_text(raw) or "").strip()
        except Exception as exc:
            raise CorruptSource("rtf", str(exc)) from exc
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class OdfLoader:
    """OpenDocument loader (ODT / ODS / ODP) on top of ``odfpy``.

    Extracts the visible text of the document body via
    ``odf.teletype.extractText``. Bind the served kind via ``kind``.
    """

    def __init__(self, *, kind: SourceKind) -> None:
        self.kind = kind

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            from odf import teletype
            from odf.opendocument import load as odf_load
        except ImportError as exc:  # pragma: no cover -- runtime dep guard
            raise CorruptSource(self.kind.value, f"odfpy is not installed: {exc}") from exc
        if isinstance(content, str):
            content = content.encode("utf-8")
        try:
            document = odf_load(io.BytesIO(content))
            body = (teletype.extractText(document.body) or "").strip()
        except Exception as exc:
            raise CorruptSource(self.kind.value, str(exc)) from exc
        if not body:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=body, order=0)], raw_text=body)


class ImageLoader:
    """Tesseract-backed OCR loader for raster images.

    Accepts PNG / JPEG / GIF / WebP (the LLM-native formats produced
    by the binary normaliser's image branch). The OCR text becomes a
    single section with ``path=[]`` -- callers that need section
    granularity for scanned multi-page documents should ingest the
    pages as separate sources or rely on the PDF intake path which
    preserves per-page boundaries.
    """

    kind = SourceKind.image

    def __init__(
        self,
        *,
        lang: str = "eng+spa",
        config: str = "--oem 3 --psm 6",
    ) -> None:
        self._lang = lang
        self._config = config

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        try:
            import pytesseract
            from PIL import Image
        except ImportError as exc:
            raise CorruptSource(
                "image",
                f"pytesseract + Pillow are required for OCR: {exc}",
            ) from exc

        if isinstance(content, str):
            content = content.encode("utf-8")

        try:
            with Image.open(io.BytesIO(content)) as img:
                if img.mode not in {"RGB", "L"}:
                    img = img.convert("RGB")
                text = pytesseract.image_to_string(img, lang=self._lang, config=self._config)
        except Exception as exc:
            raise CorruptSource("image", f"OCR failed: {exc}") from exc

        text = (text or "").strip()
        if not text:
            return LoadedDocument(raw_text="")
        return LoadedDocument(sections=[Section(path=[], body=text, order=0)], raw_text=text)


class TranscriptLoader:
    """WebVTT / SRT transcript loader.

    Stamps each cue with its timestamp and groups consecutive cues
    into a single body block per minute so chunks line up to a
    reasonable conversational unit.
    """

    kind = SourceKind.transcript

    _CUE_TIMESTAMP_RE = re.compile(
        r"(\d{1,2}:)?(\d{1,2}):(\d{2})[.,](\d{1,3})\s*-->\s*(\d{1,2}:)?(\d{1,2}):(\d{2})[.,](\d{1,3})"
    )

    def load(self, content: bytes | str, *, filename: str | None = None) -> LoadedDocument:
        if isinstance(content, bytes):
            content = content.decode("utf-8", errors="replace")
        cues: list[str] = []
        for block in content.split("\n\n"):
            block = block.strip()
            if not block:
                continue
            lines = block.splitlines()
            timestamp_line = next(
                (ln for ln in lines if self._CUE_TIMESTAMP_RE.search(ln)),
                None,
            )
            body = "\n".join(ln for ln in lines if ln and ln != timestamp_line and not ln.strip().isdigit())
            if not body.strip():
                continue
            if timestamp_line:
                cues.append(f"[{timestamp_line}] {body}")
            else:
                cues.append(body)
        if not cues:
            return LoadedDocument(raw_text="")
        body = "\n".join(cues)
        return LoadedDocument(
            sections=[Section(path=["Transcript"], body=body, order=0)],
            raw_text=body,
        )


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------


class LoaderRegistry:
    """Lookup table from :class:`SourceKind` to loader."""

    def __init__(self, fallback: SourceLoader | None = None) -> None:
        self._loaders: dict[SourceKind, SourceLoader] = {}
        self._fallback = fallback

    def register(self, loader: SourceLoader, *, kind: SourceKind | None = None) -> None:
        self._loaders[kind or loader.kind] = loader

    def set_fallback(self, loader: SourceLoader) -> None:
        self._fallback = loader

    def get(self, kind: SourceKind) -> SourceLoader | None:
        loader = self._loaders.get(kind)
        if loader is None and self._fallback is not None:
            return self._fallback
        return loader

    def __contains__(self, kind: SourceKind) -> bool:
        return kind in self._loaders or self._fallback is not None


def default_registry() -> LoaderRegistry:
    """Build the registry with every shipped loader pre-installed.

    Every format flycanon ingests has a dedicated, non-Microsoft loader.
    The fallback is a plain UTF-8 :class:`TextLoader`, so an unrecognised
    payload degrades to text extraction rather than failing. Multi-artifact
    intakes (archives, emails) arrive pre-merged as Markdown and are loaded
    by :class:`MarkdownLoader`.
    """
    registry = LoaderRegistry(fallback=TextLoader())
    registry.register(DocxLoader())
    registry.register(PdfLoader())
    registry.register(HtmlLoader())
    registry.register(MarkdownLoader())
    registry.register(TextLoader())
    registry.register(TranscriptLoader())
    registry.register(ImageLoader())
    registry.register(XlsxLoader())
    registry.register(PptxLoader())
    registry.register(CsvLoader())
    registry.register(CsvLoader(kind=SourceKind.tsv, delimiter="\t"))
    registry.register(JsonLoader())
    registry.register(XmlLoader())
    registry.register(RtfLoader())
    registry.register(OdfLoader(kind=SourceKind.odt))
    registry.register(OdfLoader(kind=SourceKind.ods))
    registry.register(OdfLoader(kind=SourceKind.odp))
    # Multi-artifact intakes (archives, emails) are merged into one Markdown
    # document by IntakeService before ingestion -- load with the Markdown
    # section detector rather than the plain-text fallback.
    registry.register(MarkdownLoader(), kind=SourceKind.archive)
    return registry


# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _markdown_to_document(text: str) -> LoadedDocument:
    """Walk Markdown text producing :class:`Section` rows by heading depth."""
    heading_re = re.compile(r"^(#{1,6})\s+(.*)$")
    sections: list[Section] = []
    heading_stack: list[str] = []
    current_path: list[str] = []
    current_body: list[str] = []
    order = 0

    def _flush() -> None:
        nonlocal order
        body = "\n".join(line for line in current_body if line.strip())
        if not body:
            return
        sections.append(Section(path=list(current_path), body=body, order=order))
        order += 1

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        match = heading_re.match(line)
        if match:
            _flush()
            current_body = []
            level = len(match.group(1))
            title = match.group(2).strip()
            while len(heading_stack) >= level:
                heading_stack.pop()
            heading_stack.append(title)
            current_path = list(heading_stack)
            continue
        current_body.append(line)
    _flush()
    return LoadedDocument(sections=sections)
